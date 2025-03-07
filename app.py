from flask import Flask, jsonify, request
from google.cloud import storage
import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
import docx
import pptx

# Initialize Flask application
app = Flask(__name__)

# Initialize AI search model and FAISS index
model = SentenceTransformer('all-MiniLM-L6-v2')

faiss_index_path = "ai_search_index.faiss"  # FAISS index file path
metadata_path = "ai_metadata.npy"  # Metadata file path (stores file paths)

try:
    index = faiss.read_index(faiss_index_path)
    file_paths = np.load(metadata_path, allow_pickle=True).tolist()  # Convert NumPy array to list
    print(f"✅ FAISS index and metadata loaded! ({len(file_paths)} files)")
except Exception as e:
    print(f"❌ Error loading FAISS index or metadata: {e}")
    index = None
    file_paths = []

# Route for homepage (testing)
@app.route('/')
def home():
    return '✅ SalesBOT API is running! Use /search to find documents.'

# 📌 Debug Route: Check if FAISS index is loaded
@app.route('/debug_index', methods=['GET'])
def debug_index():
    if index is None:
        return jsonify({"error": "FAISS index is not loaded!"}), 500
    return jsonify({"status": "FAISS index is loaded", "total_files": len(file_paths)})

# 📌 Debug Route: List all indexed files (Fix: Convert ndarray to list)
@app.route('/list_files', methods=['GET'])
def list_files():
    return jsonify({"indexed_files": list(file_paths)})  # Convert to list for JSON serialization

# 📌 Upload file route (Google Cloud Storage)
@app.route('/upload', methods=['POST'])
def upload_to_bucket():
    storage_client = storage.Client()
    bucket_name = os.environ.get('GOOGLE_STORAGE_BUCKET')
    bucket = storage_client.get_bucket(bucket_name)

    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    blob = bucket.blob(file.filename)
    blob.upload_from_file(file)

    return jsonify({"message": f"File {file.filename} uploaded successfully!"}), 200

# 📌 Function to extract text from supported file types
def extract_text_from_file(file_path):
    """Extracts text from PDFs, Word docs, and PowerPoints."""
    try:
        if file_path.endswith(".pdf"):
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file_path.endswith(".docx"):
            doc = docx.Document(file_path)
            return " ".join([para.text for para in doc.paragraphs])
        elif file_path.endswith(".pptx"):
            presentation = pptx.Presentation(file_path)
            return " ".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
        else:
            return "⚠️ Unsupported file type"
    except Exception as e:
        return f"⚠️ Error extracting text: {e}"

# 📌 Debug Route: Test text extraction on a specific file
@app.route('/extract_text', methods=['GET'])
def test_text_extraction():
    file_path = request.args.get('file_path')
    if not file_path:
        return jsonify({"error": "No file path provided"}), 400

    extracted_text = extract_text_from_file(file_path)
    return jsonify({"file_path": file_path, "extracted_text": extracted_text[:1000]})

# 📌 AI-powered search function with text extraction
def salesbot_search(query, top_k=5, file_type=None):
    if index is None:
        return [{"error": "FAISS index not loaded"}]

    query_embedding = model.encode([query], convert_to_numpy=True)
    distances, indices = index.search(query_embedding, top_k)

    results = []
    for i, idx in enumerate(indices[0]):
        if idx < len(file_paths):
            file_path = file_paths[idx]
            file_name = os.path.basename(file_path)
            drive_link = f"https://drive.google.com/open?id={file_name}"

            if file_type and not file_path.lower().endswith(file_type.lower()):
                continue

            extracted_text = extract_text_from_file(file_path)

            results.append({
                "file_name": file_name,
                "file_path": file_path,
                "google_drive_link": drive_link,
                "relevance_score": float(distances[0][i]),  # Convert float32 → float
                "document_text": extracted_text[:1000]  # Limit text for response
            })
    
    return results

# 📌 AI Search Endpoint with optional file type filtering
@app.route('/search', methods=['GET'])
def search():
    query = request.args.get('query')
    file_type = request.args.get('file_type')  # Optional file type filter

    if not query:
        return jsonify({"error": "No query provided"}), 400

    results = salesbot_search(query, file_type=file_type)
    return jsonify(results)

if __name__ == '__main__':
    app.run(debug=True, port=5000)
